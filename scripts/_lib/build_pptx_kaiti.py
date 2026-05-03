#!/usr/bin/env python3
# 把 PPT_kaiti_huibao_pandoc.md 转成精确分页的 .pptx
# 每个 ## H2 section -> 1 张幻灯片
# 支持 markdown 表格、代码块、图片

import re, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = "/Users/wuxiuxiang/project/dongmei/OptiSyn"
MD   = os.path.join(ROOT, "reports/PPT_kaiti_huibao_pandoc.md")
OUT  = os.path.join(ROOT, "reports/PPT_kaiti_huibao.pptx")

# ---------- Read & split markdown ----------
with open(MD) as f:
    md = f.read()
md = re.sub(r'^%[^\n]*\n', '', md, flags=re.MULTILINE)  # strip pandoc % metadata
sections = re.split(r'^## ', md, flags=re.MULTILINE)
sections = [s.strip() for s in sections if s.strip()]
print(f">>> {len(sections)} sections detected")

# ---------- PPTX setup (16:9, 中文字体) ----------
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

CJK_FONT = "PingFang SC"
NAVY = RGBColor(0x1A, 0x3C, 0x6E)
BLUE = RGBColor(0x2A, 0x5A, 0x9E)
RED  = RGBColor(0xD7, 0x26, 0x3D)
GRAY = RGBColor(0x44, 0x44, 0x44)
LIGHTBG = RGBColor(0xF4, 0xF6, 0xFC)
ACCENT  = RGBColor(0xF0, 0xA0, 0x20)

# ---------- Helpers ----------
def add_text(tf, text, size=14, bold=False, color=GRAY, align=None, font=CJK_FONT):
    """Add a paragraph or set first run text."""
    if not tf.paragraphs[0].runs and not tf.paragraphs[0].text:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    if align: p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def parse_inline(text):
    """Yield (text, bold) pairs splitting on **bold**."""
    parts = re.split(r'\*\*(.+?)\*\*', text)
    for i, part in enumerate(parts):
        if part:
            yield (part, i % 2 == 1)

def add_rich_paragraph(tf, text, size=14, base_color=GRAY, bold_color=NAVY, indent=0):
    """Add paragraph with **bold** inline support."""
    if tf.paragraphs[0].runs or tf.paragraphs[0].text:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    if indent: p.level = indent
    for txt, is_bold in parse_inline(text):
        run = p.add_run()
        run.text = txt
        run.font.name = CJK_FONT
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.color.rgb = bold_color if is_bold else base_color
    return p

def add_title_bar(slide, title_text):
    """Top title bar: navy bg + white text."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 prs.slide_width, Inches(0.7))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    tf = bar.text_frame; tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title_text
    r.font.name = CJK_FONT; r.font.size = Pt(22); r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

def add_table(slide, table_md, left, top, width, height, header_bold=True,
              fontsize=11):
    """Render markdown table to PPT table."""
    rows = [r.strip() for r in table_md.split('\n') if r.strip()]
    rows = [r for r in rows if not re.match(r'^\|[\s\-:|]+\|$', r)]  # drop separator
    if not rows: return None
    cells = [[c.strip() for c in re.split(r'\|', r) if c.strip() != '' or r.startswith('|')]
             for r in rows]
    cells = [[c for c in row if c != ''] for row in cells]
    nrow, ncol = len(cells), max(len(r) for r in cells) if cells else 0
    if nrow == 0 or ncol == 0: return None
    tbl = slide.shapes.add_table(nrow, ncol, left, top, width, height).table
    for i, row in enumerate(cells):
        for j in range(ncol):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHTBG if i == 0 else RGBColor(0xFF,0xFF,0xFF)
            tf = cell.text_frame
            tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
            text = row[j] if j < len(row) else ''
            tf.text = ''
            for k, (txt, b) in enumerate(parse_inline(text)):
                if k == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph() if False else None
                run = (tf.paragraphs[0]).add_run() if k == 0 else \
                      (tf.paragraphs[0]).add_run()
                run.text = txt
                run.font.name = CJK_FONT
                run.font.size = Pt(fontsize)
                run.font.bold = (i == 0 and header_bold) or b
                if i == 0: run.font.color.rgb = NAVY
                else: run.font.color.rgb = GRAY
    return tbl

def parse_section(section):
    """Returns (title, blocks) where blocks is a list of dicts:
       {kind: 'text'|'bullet'|'table'|'code'|'image', content: ...}"""
    lines = section.split('\n')
    title = lines[0].strip()
    body = '\n'.join(lines[1:]).strip()
    blocks = []
    i = 0
    body_lines = body.split('\n')
    while i < len(body_lines):
        line = body_lines[i]
        # code block
        if line.startswith('```'):
            j = i + 1
            code = []
            while j < len(body_lines) and not body_lines[j].startswith('```'):
                code.append(body_lines[j])
                j += 1
            blocks.append({'kind': 'code', 'content': '\n'.join(code)})
            i = j + 1
            continue
        # table
        if line.startswith('|'):
            j = i
            tab = []
            while j < len(body_lines) and body_lines[j].startswith('|'):
                tab.append(body_lines[j])
                j += 1
            blocks.append({'kind': 'table', 'content': '\n'.join(tab)})
            i = j
            continue
        # image
        m = re.match(r'!\[\]?\(([^)]+)\)', line)
        if m:
            blocks.append({'kind': 'image', 'content': m.group(1)})
            i += 1
            continue
        # bullet
        if re.match(r'^[\-\*]\s+', line):
            j = i
            bs = []
            while j < len(body_lines) and re.match(r'^[\-\*]\s+', body_lines[j]):
                bs.append(re.sub(r'^[\-\*]\s+', '', body_lines[j]))
                j += 1
            blocks.append({'kind': 'bullet', 'content': bs})
            i = j
            continue
        # numbered list
        if re.match(r'^\d+\.\s+', line):
            j = i
            bs = []
            while j < len(body_lines) and re.match(r'^\d+\.\s+', body_lines[j]):
                bs.append(re.sub(r'^\d+\.\s+', '', body_lines[j]))
                j += 1
            blocks.append({'kind': 'numbered', 'content': bs})
            i = j
            continue
        # text paragraph (collect non-empty consecutive lines)
        if line.strip():
            blocks.append({'kind': 'text', 'content': line.strip()})
        i += 1
    return title, blocks

# ---------- Build slides ----------
for idx, sec in enumerate(sections):
    title, blocks = parse_section(sec)
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    # 第一张是封面
    if idx == 0 or title.startswith("封面"):
        # cover
        cover_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                           prs.slide_width, prs.slide_height)
        cover_bar.fill.solid(); cover_bar.fill.fore_color.rgb = NAVY
        cover_bar.line.fill.background()
        # main title
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(2.0),
                                       Inches(12.1), Inches(2.2))
        tf = tx.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = '靶向"外周肺-中枢孤束核"神经免疫互作轴的'
        r.font.name = CJK_FONT; r.font.size = Pt(32); r.font.bold = True
        r.font.color.rgb = RGBColor(255, 255, 255)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = '流感后咳嗽机制及中药干预深度报告'
        r2.font.name = CJK_FONT; r2.font.size = Pt(32); r2.font.bold = True
        r2.font.color.rgb = RGBColor(255, 255, 255)
        # subtitle
        tx2 = slide.shapes.add_textbox(Inches(0.6), Inches(4.5),
                                        Inches(12.1), Inches(1))
        tf2 = tx2.text_frame
        p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = "项目开题汇报"
        r.font.name = CJK_FONT; r.font.size = Pt(22)
        r.font.color.rgb = ACCENT
        # meta
        tx3 = slide.shapes.add_textbox(Inches(0.6), Inches(5.5),
                                        Inches(12.1), Inches(1.5))
        tf3 = tx3.text_frame
        for line in ["汇报人: xiuxiang  ·  2026-05-03",
                     "数据: GSE42639 / GSE31022 / GSE161878 / GSE296065 / GSE268741",
                     "         + CellChatDB.mouse"]:
            p = tf3.add_paragraph() if tf3.paragraphs[0].text else tf3.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = line
            r.font.name = CJK_FONT; r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(220, 220, 220)
        continue

    # Regular slide: title bar + body
    add_title_bar(slide, title)

    # body region
    bx = Inches(0.4); by = Inches(0.9)
    bw = prs.slide_width - Inches(0.8)
    bh = prs.slide_height - Inches(1.2)

    cur_y = by
    for blk in blocks:
        kind = blk['kind']
        if kind == 'text':
            tx = slide.shapes.add_textbox(bx, cur_y, bw, Inches(0.5))
            tf = tx.text_frame; tf.word_wrap = True
            tf.margin_top = Inches(0.05)
            add_rich_paragraph(tf, blk['content'], size=14)
            cur_y += Inches(0.5)
        elif kind == 'bullet':
            n = len(blk['content'])
            tx = slide.shapes.add_textbox(bx, cur_y, bw, Inches(0.35 * n + 0.15))
            tf = tx.text_frame; tf.word_wrap = True
            tf.margin_top = Inches(0.05)
            for k, b in enumerate(blk['content']):
                if k == 0 and not (tf.paragraphs[0].runs or tf.paragraphs[0].text):
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.level = 0
                # bullet glyph
                rb = p.add_run(); rb.text = "▸ "
                rb.font.name = CJK_FONT; rb.font.size = Pt(13)
                rb.font.color.rgb = ACCENT
                for txt, is_bold in parse_inline(b):
                    run = p.add_run(); run.text = txt
                    run.font.name = CJK_FONT; run.font.size = Pt(13)
                    run.font.bold = is_bold
                    run.font.color.rgb = NAVY if is_bold else GRAY
            cur_y += Inches(0.35 * n + 0.20)
        elif kind == 'numbered':
            n = len(blk['content'])
            tx = slide.shapes.add_textbox(bx, cur_y, bw, Inches(0.35 * n + 0.15))
            tf = tx.text_frame; tf.word_wrap = True
            tf.margin_top = Inches(0.05)
            for k, b in enumerate(blk['content']):
                if k == 0 and not (tf.paragraphs[0].runs or tf.paragraphs[0].text):
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                rb = p.add_run(); rb.text = f"{k+1}. "
                rb.font.name = CJK_FONT; rb.font.size = Pt(13); rb.font.bold = True
                rb.font.color.rgb = BLUE
                for txt, is_bold in parse_inline(b):
                    run = p.add_run(); run.text = txt
                    run.font.name = CJK_FONT; run.font.size = Pt(13)
                    run.font.bold = is_bold
                    run.font.color.rgb = NAVY if is_bold else GRAY
            cur_y += Inches(0.35 * n + 0.20)
        elif kind == 'table':
            tab_lines = [l for l in blk['content'].split('\n')
                         if not re.match(r'^\|[\s\-:|]+\|$', l)]
            n_rows = len(tab_lines)
            est_h = Inches(0.35 * n_rows + 0.1)
            add_table(slide, blk['content'], bx, cur_y, bw, est_h,
                      fontsize=11)
            cur_y += est_h + Inches(0.10)
        elif kind == 'code':
            n_lines = len(blk['content'].split('\n'))
            est_h = Inches(0.20 * n_lines + 0.20)
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, cur_y,
                                          bw, est_h)
            box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF4,0xF4,0xF6)
            box.line.color.rgb = RGBColor(0xCC,0xCC,0xDD)
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.05)
            for k, line in enumerate(blk['content'].split('\n')):
                if k == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                r = p.add_run(); r.text = line if line else ' '
                r.font.name = "Menlo"; r.font.size = Pt(10)
                r.font.color.rgb = RGBColor(0x33, 0x33, 0x55)
            cur_y += est_h + Inches(0.10)
        elif kind == 'image':
            img_path = blk['content']
            if not os.path.isabs(img_path):
                img_path = os.path.join(os.path.dirname(MD), img_path)
                img_path = os.path.abspath(img_path)
            if os.path.exists(img_path):
                # resize to fit
                avail_h = prs.slide_height - cur_y - Inches(0.4)
                img = slide.shapes.add_picture(img_path, bx, cur_y,
                                                width=bw, height=avail_h)
                cur_y += avail_h
            else:
                print(f"  WARN: image not found: {img_path}")

# ---------- Save ----------
prs.save(OUT)
print(f">>> 输出: {OUT}")
print(f">>> 总幻灯片数: {len(prs.slides)}")
